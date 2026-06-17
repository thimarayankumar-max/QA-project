from pypdf import PdfReader
from docx import Document
from pptx import Presentation

TXT_LIMIT = 10 * 1024 * 1024

def extract_text_from_file(file):
    filename = file.filename.lower()

    try:
        if filename.endswith(".pdf"):
            reader = PdfReader(file)
            text = ""
            for page in reader.pages[:50]:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + " "
            return text if text.strip() else "Error: No text found in PDF."

        if filename.endswith(".docx"):
            doc = Document(file)
            text = ""
            for para in doc.paragraphs:
                if para.text.strip():
                    text += para.text + " "
            return text if text.strip() else "Error: No text found in DOCX."

        if filename.endswith(".pptx"):
            ppt = Presentation(file)
            text = ""
            for slide in ppt.slides[:50]:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text + " "
            return text if text.strip() else "Error: No text found in PPTX."

        if filename.endswith(".txt"):
            data = file.read()
            if len(data) > TXT_LIMIT:
                return "Error: TXT file size exceeds 10 MB limit."
            return data.decode("utf-8", errors="ignore")

        return "Unsupported file type"

    except Exception as e:
        return "Error: " + str(e)